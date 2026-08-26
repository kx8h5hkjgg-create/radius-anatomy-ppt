from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with custom settings
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(192, 0, 60)  # Pink/Maroon
SECONDARY_COLOR = RGBColor(100, 149, 237)  # Cornflower Blue
ACCENT_COLOR = RGBColor(220, 20, 60)  # Crimson
TEXT_COLOR = RGBColor(40, 40, 40)  # Dark Gray
LIGHT_BG = RGBColor(245, 245, 250)  # Light background

def add_animations_to_shape(shape, animation_type="appear"):
    """Add animations to shapes"""
    pass  # python-pptx has limited animation support

def add_title_slide(prs, title, subtitle):
    """Add animated title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Add decorative shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(6.5),
        Inches(10), Inches(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.color.rgb = ACCENT_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list, add_icon=False):
    """Add content slide with bullet points and styling"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Add title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.9)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Add accent line
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.9),
        Inches(10), Inches(0.08)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT_COLOR
    accent_line.line.color.rgb = ACCENT_COLOR
    
    # Title text
    title_frame = title_shape.text_frame
    title_frame.margin_bottom = Inches(0.05)
    title_frame.margin_top = Inches(0.05)
    title_frame.margin_left = Inches(0.3)
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(10)
        p.space_after = Pt(10)
        p.line_spacing = 1.35
    
    return slide

def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """Add two-column layout slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Add title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.9)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Add accent line
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.9),
        Inches(10), Inches(0.08)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT_COLOR
    accent_line.line.color.rgb = ACCENT_COLOR
    
    title_frame = title_shape.text_frame
    title_frame.margin_bottom = Inches(0.05)
    title_frame.margin_top = Inches(0.05)
    title_frame.margin_left = Inches(0.3)
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left column box
    left_box_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.1),
        Inches(4.6), Inches(5.8)
    )
    left_box_shape.fill.solid()
    left_box_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    left_box_shape.line.color.rgb = SECONDARY_COLOR
    left_box_shape.line.width = Pt(2)
    
    # Left title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(0.4))
    left_title_frame = left_title_box.text_frame
    p = left_title_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    # Left content
    left_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(5))
    left_frame = left_text_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.line_spacing = 1.2
    
    # Right column box
    right_box_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.1), Inches(1.1),
        Inches(4.6), Inches(5.8)
    )
    right_box_shape.fill.solid()
    right_box_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    right_box_shape.line.color.rgb = SECONDARY_COLOR
    right_box_shape.line.width = Pt(2)
    
    # Right title
    right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.3), Inches(4.2), Inches(0.4))
    right_title_frame = right_title_box.text_frame
    p = right_title_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    # Right content
    right_text_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(4.2), Inches(5))
    right_frame = right_text_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.line_spacing = 1.2
    
    return slide

def add_highlight_slide(prs, title, main_text, highlights):
    """Add slide with highlighted key points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG
    
    # Add title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.9)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Add accent line
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.9),
        Inches(10), Inches(0.08)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = ACCENT_COLOR
    accent_line.line.color.rgb = ACCENT_COLOR
    
    title_frame = title_shape.text_frame
    title_frame.margin_bottom = Inches(0.05)
    title_frame.margin_top = Inches(0.05)
    title_frame.margin_left = Inches(0.3)
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Main text
    main_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(1.2))
    main_frame = main_box.text_frame
    main_frame.word_wrap = True
    p = main_frame.paragraphs[0]
    p.text = main_text
    p.font.size = Pt(22)
    p.font.italic = True
    p.font.color.rgb = SECONDARY_COLOR
    
    # Highlight boxes
    for i, highlight in enumerate(highlights):
        y_pos = 2.6 + (i * 1.3)
        
        # Highlight box
        highlight_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(y_pos),
            Inches(8.6), Inches(1.1)
        )
        highlight_shape.fill.solid()
        highlight_shape.fill.fore_color.rgb = RGBColor(255, 250, 200)
        highlight_shape.line.color.rgb = ACCENT_COLOR
        highlight_shape.line.width = Pt(2)
        
        # Highlight text
        highlight_box = slide.shapes.add_textbox(Inches(0.9), Inches(y_pos + 0.15), Inches(8.2), Inches(0.8))
        highlight_frame = highlight_box.text_frame
        highlight_frame.word_wrap = True
        p = highlight_frame.paragraphs[0]
        p.text = highlight
        p.font.size = Pt(19)
        p.font.bold = True
        p.font.color.rgb = ACCENT_COLOR
        p.line_spacing = 1.2
    
    return slide

# ==================== SLIDE CONTENT ====================

# Slide 1: Title Slide
add_title_slide(prs, "THE RADIUS BONE", "Comprehensive Anatomy Guide")

# Slide 2: Introduction
add_content_slide(prs, "Introduction to the Radius", [
    "• The LATERAL bone of the forearm",
    "• Pre-axial bone of lower limb",
    "• Corresponds with the tibia",
    "• Forms the lateral (thumb) side of forearm",
    "• Works with ulna for pronation and supination",
    "• Key articulation with humerus (elbow joint)"
])

# Slide 3: General Overview
add_content_slide(prs, "Basic Structure", [
    "• THREE main regions:",
    "    ◆ Upper (proximal) end",
    "    ◆ Shaft (diaphysis)",
    "    ◆ Lower (distal) end",
    "",
    "• KEY FUNCTIONS:",
    "    ◆ Articulates with humerus at elbow",
    "    ◆ Rotates around ulna",
    "    ◆ Supports wrist and hand"
])

# Slide 4: Anatomical Position
add_content_slide(prs, "Anatomical Position & Orientation", [
    "• When forearm is in SUPINATION:",
    "    ◆ Radius is lateral (thumb side)",
    "    ◆ Ulna is medial",
    "",
    "• When forearm is in PRONATION:",
    "    ◆ Radius crosses over ulna",
    "    ◆ Rotation occurs at radio-ulnar joints",
    "",
    "• Understanding position crucial for anatomy"
])

# Slide 5: Upper End - Overview
add_content_slide(prs, "Upper End (Proximal End) - Overview", [
    "• CONSISTS of THREE distinct parts:",
    "    1. Head (Caput radii)",
    "    2. Neck (Collum)",
    "    3. Radial tuberosity",
    "",
    "• Articulates with humerus and ulna",
    "• Supports pronation and supination",
    "• Contains ligamentous attachments"
])

# Slide 6: The Radial Head - Basic Features
add_content_slide(prs, "The Radial Head (Caput) - Basic Features", [
    "• DISC-LIKE structure",
    "• CONCAVE upper surface",
    "• Peripheral margin BROADER on medial side",
    "",
    "• Upper surface articulates with:",
    "    ◆ Capitulum of humerus",
    "    ◆ Forms humero-radial part of elbow joint",
    "",
    "• Posterior surface palpable during",
    "    pronation and supination movements"
])

# Slide 7: Radial Head - Detailed Articular Surfaces
add_content_slide(prs, "Radial Head - Articular Surfaces", [
    "• UPPER SURFACE (concave disc):",
    "    - Articulates with capitulum of humerus",
    "",
    "• PERIPHERAL MARGIN:",
    "    - Medially: Articulates with radial notch of ulna",
    "    - Rest encircled by ANNULAR LIGAMENT",
    "",
    "• Forms SUPERIOR RADIO-ULNAR JOINT",
    "• Palpable through lateral depressed area"
])

# Slide 8: Neck of Radius - Structure
add_content_slide(prs, "The Neck (Collum) - Structure", [
    "• CONSTRICTED area below the head",
    "• Encircled by LOWER PART OF ANNULAR LIGAMENT",
    "",
    "• SEPARATED by:",
    "    ◆ Synovial protrusion of superior radio-ulnar joint",
    "",
    "• SUPPORTED by:",
    "    ◆ QUADRATE LIGAMENT",
    "    ◆ Derived from interlacement of annular ligament"
])

# Slide 9: Neck of Radius - Clinical Importance
add_content_slide(prs, "The Neck - Clinical Relevance", [
    "• Narrow constriction increases fracture risk",
    "• Important landmark for anatomical identification",
    "• Ligamentous support prevents excessive motion",
    "",
    "• QUADRATE LIGAMENT functions:",
    "    ◆ Stabilizes the neck",
    "    ◆ Supports radio-ulnar joint",
    "    ◆ Limits excessive rotation"
])

# Slide 10: Radial Tuberosity - Part 1
add_content_slide(prs, "Radial Tuberosity - Overview", [
    "• Located on MEDIAL side of lower neck",
    "• ROUGH elevation on bone surface",
    "",
    "• DIVIDED into TWO parts:",
    "    1. POSTERIOR part - Rough surface",
    "    2. ANTERIOR part - Smooth surface",
    "",
    "• Important muscular attachment point",
    "• Palpable during physical examination"
])

# Slide 11: Radial Tuberosity - Detailed Features
add_content_slide(prs, "Radial Tuberosity - Muscular Attachments", [
    "• POSTERIOR PART (Rough):",
    "    ◆ Receives TENDON OF BICEPS BRACHII",
    "    ◆ Twisted manner of insertion",
    "    ◆ Powerful attachment for elbow flexion",
    "",
    "• ANTERIOR PART (Smooth):",
    "    ◆ Separated from biceps by BURSA",
    "    ◆ Bursa allows smooth supination",
    "    ◆ Prevents friction with biceps tendon"
])

# Slide 12: Radial Tuberosity - Cord Attachment
add_content_slide(prs, "Radial Tuberosity - Oblique Cord", [
    "• LOWER END of tuberosity:",
    "    ◆ Gives attachment to OBLIQUE CORD",
    "",
    "• OBLIQUE CORD:",
    "    ◆ Runs from radius to ulna",
    "    ◆ Complements interosseous membrane",
    "    ◆ Provides additional stability",
    "",
    "• Functional importance in forearm support"
])

# Slide 13: Shaft - General Features
add_content_slide(prs, "The Shaft (Diaphysis) - Overview", [
    "• TRIANGULAR cross-section in middle third",
    "",
    "• THREE BORDERS:",
    "    1. Anterior border (oblique to vertical)",
    "    2. Posterior border (prominent in middle)",
    "    3. Interosseous border (medial)",
    "",
    "• THREE SURFACES:",
    "    1. Anterior surface (concave)",
    "    2. Posterior surface",
    "    3. Lateral surface (convex)"
])

# Slide 14: Anterior Border - Upper Part
add_content_slide(prs, "Anterior Border - Upper Portion", [
    "• OBLIQUE orientation in upper part",
    "• Slopes DOWNWARD and LATERALLY",
    "• Originates from radial tuberosity",
    "",
    "• MUSCULAR ATTACHMENT:",
    "    ◆ Upper oblique portion gives origin to",
    "    ◆ RADIAL HEAD OF FLEXOR DIGITORUM SUPERFICIALIS",
    "",
    "• Important for finger flexion function",
    "• Provides broad attachment area"
])

# Slide 15: Anterior Border - Lower Part
add_content_slide(prs, "Anterior Border - Lower Portion", [
    "• Changes from OBLIQUE to VERTICAL",
    "• Extends as PROMINENT RIDGE",
    "• Located in distal fourth of shaft",
    "",
    "• CONTINUOUS with:",
    "    ◆ Anterior border of styloid process",
    "",
    "• ATTACHMENT:",
    "    ◆ Lateral end of EXTENSOR RETINACULUM",
    "    ◆ Important for wrist extension support"
])

# Slide 16: Posterior Border
add_content_slide(prs, "Posterior Border - Characteristics", [
    "• PROMINENT in the MIDDLE THIRD",
    "• ILL-DEFINED in upper and lower thirds",
    "",
    "• WHEN TRACED UPWARD:",
    "    ◆ Slopes upward and medially",
    "    ◆ Reaches postero-inferior part of tuberosity",
    "",
    "• ANATOMICAL SIGNIFICANCE:",
    "    ◆ Landmark for side determination",
    "    ◆ Provides muscular attachment area"
])

# Slide 17: Interosseous Border - Primary Features
add_content_slide(prs, "Interosseous (Medial) Border - Overview", [
    "• SHARP border",
    "• Extends from BELOW RADIAL TUBEROSITY to lower end",
    "",
    "• Connected to ulna by:",
    "    ◆ INTEROSSEOUS MEMBRANE (strong)",
    "    ◆ OBLIQUE CORD (upper part)",
    "",
    "• Provides extensive surface area",
    "• Critical for forearm stability"
])

# Slide 18: Interosseous Border - Vascular Features
add_content_slide(prs, "Interosseous Border - Vascular Anatomy", [
    "• UPPER MARGIN: FREE (gap present)",
    "",
    "• GAP between radius and oblique cord transmits:",
    "    ◆ POSTERIOR INTEROSSEOUS ARTERY (BLOOD VESSEL)",
    "    ◆ NOT the posterior interosseous nerve",
    "",
    "• LOWER MARGIN:",
    "    ◆ Continuous with inferior radio-ulnar joint",
    "    ◆ Connected to capsule ligament"
])

# Slide 19: Interosseous Membrane - Functional Role
add_content_slide(prs, "Interosseous Membrane - Structure & Function", [
    "• STRONG fibrous membrane between radius & ulna",
    "",
    "• UPPER MARGIN: FREE (allows vessel passage)",
    "",
    "• LOWER MARGIN:",
    "    ◆ Continuous with capsule of inferior radio-ulnar joint",
    "    ◆ Separated anteriorly by synovial pouch",
    "    ◆ Recessus sacciformis",
    "",
    "• Provides strength and stability to forearm"
])

# Slide 20: Anterior Surface of Shaft
add_content_slide(prs, "Anterior Surface of Shaft", [
    "• Located between ANTERIOR and INTEROSSEOUS borders",
    "• GENTLY CONCAVE contour",
    "",
    "• NUTRIENT FORAMEN:",
    "    ◆ Points TOWARDS the elbow",
    "    ◆ Located near MIDDLE of shaft",
    "    ◆ Allows vascular entry",
    "",
    "• Provides smooth gliding surface",
    "• Important muscular attachment area"
])

# Slide 21: Posterior Surface - Upper Two-Thirds
add_content_slide(prs, "Posterior Surface - Upper Portion", [
    "• Between POSTERIOR and INTEROSSEOUS borders",
    "",
    "• MUSCULAR ATTACHMENTS:",
    "    ◆ Upper two-thirds",
    "    ◆ Anterior surface with triangular medial area",
    "    ◆ In front of interosseous membrane",
    "",
    "• PRONATOR QUADRATUS:",
    "    ◆ Inserts on anterior surface",
    "    ◆ Critical for pronation movement"
])

# Slide 22: Posterior Surface - Lower Portion
add_content_slide(prs, "Posterior Surface - Lower Portion", [
    "• Lower one-fourth of posterior surface",
    "",
    "• GIVES ORIGIN TO:",
    "    ◆ FLEXOR POLLICIS LONGUS",
    "    ◆ Important for thumb flexion",
    "",
    "• Attachment to pronator quadratus",
    "",
    "• Clinical importance for hand function"
])

# Slide 23: Lateral Surface - Contour
add_content_slide(prs, "Lateral Surface - General Features", [
    "• GENTLY CONVEX surface",
    "• Summit (highest point) at MIDDLE of shaft",
    "",
    "• ROUGHENED impressions on surface",
    "",
    "• Important muscular attachment region",
    "",
    "• Palpable during physical examination",
    "• Reference for anatomical orientation"
])

# Slide 24: Lateral Surface - Muscular Attachments
add_content_slide(prs, "Lateral Surface - Muscle Attachments", [
    "• UPPER PART:",
    "    ◆ Encroached impression for PRONATOR TERES",
    "    ◆ Provides origin for this muscle",
    "",
    "• SUPINATOR MUSCLE:",
    "    ◆ Receives insertion on lateral surface",
    "    ◆ Deep fibers pass around radius",
    "",
    "• DEEP RADIAL NERVE:",
    "    ◆ Passes through supinator",
    "    ◆ Can be compressed here"
])

# Slide 25: Lower End - General Overview
add_title_slide(prs, "LOWER END OF RADIUS", "Distal Anatomy & Articular Surfaces")

# Slide 26: Lower End - Four Surfaces
add_content_slide(prs, "Lower End (Distal End) - Four Surfaces", [
    "• WIDEST PART of the entire bone",
    "",
    "• Presents FOUR surfaces:",
    "    1. LATERAL - Styloid process",
    "    2. MEDIAL - Ulnar notch",
    "    3. ANTERIOR - Palmar aspect",
    "    4. POSTERIOR - Dorsal aspect with grooves",
    "",
    "• INFERIOR CARPAL ARTICULAR SURFACE",
    "• Critical for wrist function"
])

# Slide 27: Lateral Surface - Styloid Process
add_content_slide(prs, "Lateral Surface - Styloid Process", [
    "• ROUGH surface appearance",
    "• Projects DOWNWARD as STYLOID PROCESS",
    "• Extends BEYOND styloid of ulna",
    "",
    "• STYLOID PROCESS TIP:",
    "    ◆ Attachment to LATERAL CARPAL LIGAMENT",
    "",
    "• PROXIMAL TO STYLOID:",
    "    ◆ Insertion of BRACHIORADIALIS",
    "    ◆ Crossed by ABDUCTOR POLLICIS LONGUS",
    "    ◆ Crossed by EXTENSOR POLLICIS BREVIS"
])

# Slide 28: Medial Surface - Ulnar Notch
add_content_slide(prs, "Medial Surface - Ulnar Notch Features", [
    "• Located distal to triangular area",
    "",
    "• Presents ULNAR NOTCH:",
    "    ◆ Articulates with HEAD OF ULNA",
    "    ◆ Forms INFERIOR RADIO-ULNAR JOINT",
    "",
    "• JUNCTION POINT:",
    "    ◆ Between ulnar notch and carpal surface",
    "    ◆ Attachment to TRIANGULAR ARTICULAR DISC",
    "",
    "• Ulna is EXCLUDED from wrist joint formation"
])

# Slide 29: Medial Surface - Articular Disc
add_content_slide(prs, "Medial Surface - Triangular Articular Disc", [
    "• TRIANGULAR FIBROUS DISC of inferior radio-ulnar joint",
    "",
    "• APEX fixed to depression between:",
    "    ◆ Inferior articular surface of head of ulna",
    "    ◆ Styloid process of ulna",
    "",
    "• PROVIDES:",
    "    ◆ Stability to radio-ulnar joint",
    "    ◆ Separation of radio-ulnar and wrist joints",
    "    ◆ Support for carpal articulation"
])

# Slide 30: Anterior Surface - Ridge & Ligament
add_content_slide(prs, "Anterior Surface - Palmar Features", [
    "• Represented by THICK PROMINENT RIDGE",
    "• PALPABLE through overlying tendons",
    "",
    "• ATTACHMENT to:",
    "    ◆ PALMAR RADIO-CARPAL LIGAMENT",
    "    ◆ Supports wrist joint",
    "",
    "• RADIAL ARTERY PULSATION:",
    "    ◆ Felt against this surface",
    "    ◆ Distal to pronator quadratus",
    "    ◆ Important clinical landmark"
])

# Slide 31: Posterior Surface - Dorsal Tubercle
add_content_slide(prs, "Posterior Surface - Dorsal Tubercle (Lister)", [
    "• PALPABLE DORSAL TUBERCLE",
    "• Also known as TUBERCLE OF LISTER",
    "",
    "• Located on posterior surface of lower end",
    "",
    "• SERVES AS PULLEY for:",
    "    ◆ EXTENSOR POLLICIS LONGUS tendon",
    "    ◆ Redirects tendon path to thumb",
    "",
    "• Important anatomical landmark",
    "• Easily palpable during physical exam"
])

# Slide 32: Posterior Surface - Grooves
add_content_slide(prs, "Posterior Surface - Three Longitudinal Grooves", [
    "• Displays THREE GROOVES for extensor tendons:",
    "",
    "• LATERAL GROOVE (wide):",
    "    ◆ Lateral to dorsal tubercle",
    "    ◆ Houses EXTENSOR CARPI RADIALIS LONGUS",
    "    ◆ Houses EXTENSOR CARPI RADIALIS BREVIS",
    "    ◆ Ridge intervenes between them",
    "",
    "• Grooves serve as fibro-osseous tunnels"
])

# Slide 33: Posterior Surface - Medial Grooves
add_content_slide(prs, "Posterior Surface - Medial Grooves", [
    "• GROOVE JUST MEDIAL TO DORSAL TUBERCLE:",
    "    ◆ CONSPICUOUS and deep",
    "    ◆ Transmits EXTENSOR POLLICIS LONGUS tendon",
    "    ◆ Uses tubercle as pulley",
    "",
    "• GROOVE MORE MEDIALLY (shallow):",
    "    ◆ Occupied by EXTENSOR DIGITORUM",
    "    ◆ More deeply by EXTENSOR INDICIS",
    "    ◆ POSTERIOR INTEROSSEOUS NERVE passes here"
])

# Slide 34: Extensor Retinaculum
add_content_slide(prs, "Extensor Retinaculum - Attachment & Function", [
    "• Strong fibrous band on dorsal wrist",
    "",
    "• ALL extensor tendons pass BENEATH it",
    "",
    "• GIVES SLIP OF ATTACHMENT to:",
    "    ◆ Dorsal tubercle (forms additional pulley)",
    "",
    "• STRUCTURES SECURED:",
    "    ◆ Multiple extensor tendons",
    "    ◆ Dorsal branch of radial nerve",
    "    ◆ Posterior interosseous nerve",
    "",
    "• Forms fibro-osseous tunnels for smooth movement"
])

# Slide 35: Inferior Carpal Articular Surface - Structure
add_content_slide(prs, "Inferior Carpal Articular Surface - Basic Features", [
    "• CONCAVE overall contour",
    "• Articulates with carpal bones",
    "",
    "• SUBDIVIDED by RIDGE into TWO areas:",
    "    1. LATERAL triangular area",
    "    2. MEDIAL quadrilateral area",
    "",
    "• Surface adapted for wrist motion",
    "• Allows complex hand movements"
])

# Slide 36: Carpal Articular Surface - Specific Articulations
add_content_slide(prs, "Carpal Articular Surface - Bone Articulations", [
    "• LATERAL TRIANGULAR AREA:",
    "    ◆ Articulates with SCAPHOID bone",
    "",
    "• MEDIAL QUADRILATERAL AREA:",
    "    ◆ Articulates with LUNATE bone",
    "",
    "• TOGETHER form:",
    "    ◆ RADIO-CARPAL JOINT (Wrist joint)",
    "    ◆ Most important joint of the hand",
    "",
    "• Allows flexion, extension, radial/ulnar deviation"
])

# Slide 37: Carpal Articular Surface - Attachments
add_content_slide(prs, "Carpal Articular Surface - Ligamentous Attachments", [
    "• FIBROUS CAPSULE:",
    "    ◆ Attached along entire periphery",
    "    ◆ At ulnar notch and margins",
    "",
    "• ANTERIOR ATTACHMENT:",
    "    ◆ Palmar radio-carpal ligaments",
    "",
    "• POSTERIOR ATTACHMENT:",
    "    ◆ Dorsal radio-carpal ligaments",
    "",
    "• Creates secure articulation for wrist"
])

# Slide 38: Inferior Radio-Ulnar Joint
add_content_slide(prs, "Inferior Radio-Ulnar Joint - Key Features", [
    "• Articulation between radius and ulna at wrist",
    "",
    "• JOINT SURFACES:",
    "    ◆ Ulnar notch of radius",
    "    ◆ Head of ulna",
    "",
    "• SEPARATED from wrist joint by:",
    "    ◆ TRIANGULAR ARTICULAR DISC",
    "",
    "• ALLOWS:",
    "    ◆ Pronation movement",
    "    ◆ Supination movement",
    "    ◆ Minimal gliding motion"
])

# Slide 39: Important Anatomical Relationships
add_two_column_slide(prs, "Upper End Articulations & Relationships", 
    "PROXIMALLY",
    [
        "• Head of radius articulates with CAPITULUM OF HUMERUS",
        "• Forms HUMERO-RADIAL part of ELBOW JOINT",
        "• Encircled by ANNULAR LIGAMENT",
        "• Supported by RADIAL COLLATERAL LIGAMENT",
        "• Allows flexion/extension at elbow",
        "• Peripheral margin articulates with ULNAR NOTCH"
    ],
    "SUPPORTING STRUCTURES",
    [
        "• Annular ligament encircles radial head",
        "• Quadrate ligament supports neck",
        "• Interosseous membrane connects to ulna",
        "• Oblique cord provides additional support",
        "• Radial collateral ligament reinforces joint",
        "• Provides rotational stability"
    ]
)

# Slide 40: Important Anatomical Relationships II
add_two_column_slide(prs, "Lower End Articulations & Relationships",
    "DISTALLY (WRIST JOINT)",
    [
        "• Radius articulates with SCAPHOID",
        "• Radius articulates with LUNATE",
        "• Excluded: ULNA from wrist joint",
        "• Ulna connected via triangular disc",
        "• Forms primary wrist joint",
        "• Allows complex wrist movements"
    ],
    "DISTAL RADIO-ULNAR JOINT",
    [
        "• Radius articulates with ulnar head",
        "• Separated from wrist joint",
        "• Triangular disc provides separation",
        "• Allows pronation/supination",
        "• Gliding movement only",
        "• Synovial membrane involvement"
    ]
)

# Slide 41: Side Determination - Method
add_highlight_slide(prs, "How to Determine the Side of Radius",
    "Place the disc-like head of the radius ABOVE",
    [
        "✓ Gentle CONCAVITY of shaft in FRONT → Lateral side",
        "✓ STYLOID PROCESS projecting LATERALLY → Lateral side",
        "✓ Anterior border vertical ridge → Distal aspect",
        "✓ Rough anterior border → Upper end identification"
    ]
)

# Slide 42: Clinical Significance - Fractures
add_content_slide(prs, "Clinical Significance - Common Fractures", [
    "• COLLES' FRACTURE:",
    "    ◆ Most common wrist fracture",
    "    ◆ Distal radius fracture with dorsal displacement",
    "",
    "• SMITH'S FRACTURE:",
    "    ◆ Reverse of Colles' fracture",
    "    ◆ Ventral/palmar displacement",
    "",
    "• MONTEGGIA FRACTURE:",
    "    ◆ Radius fracture + ulnar dislocation",
    "    ◆ Higher energy injury"
])

# Slide 43: Clinical Significance - Nerve Involvement
add_content_slide(prs, "Clinical Significance - Nerve Compression", [
    "• POSTERIOR INTEROSSEOUS NERVE:",
    "    ◆ Compressed by supinator muscle",
    "    ◆ Results in wrist drop",
    "    ◆ Loss of finger/thumb extension",
    "",
    "• ANTERIOR INTEROSSEOUS NERVE:",
    "    ◆ Pronator syndrome involvement",
    "    ◆ Loss of thumb IP flexion",
    "",
    "• RADIAL NERVE:",
    "    ◆ Can be injured in proximal fractures"
])

# Slide 44: Clinical Significance - Vascular
add_content_slide(prs, "Clinical Significance - Vascular Assessment", [
    "• RADIAL ARTERY PULSATION:",
    "    ◆ Palpated on anterior surface near wrist",
    "    ◆ Lateral to flexor carpi radialis tendon",
    "    ◆ Important vital sign assessment",
    "",
    "• POSTERIOR INTEROSSEOUS ARTERY:",
    "    ◆ Passes through gap in interosseous border",
    "    ◆ Supplies posterior forearm compartment",
    "",
    "• Preserved in most fracture management"
])

# Slide 45: Muscular Attachments - Summary Table
add_two_column_slide(prs, "Complete Muscular Attachments Summary",
    "PROXIMAL ATTACHMENTS",
    [
        "✓ Biceps brachii - Tuberosity (powerful)",
        "✓ Supinator muscle - Upper lateral shaft",
        "✓ Flexor digitorum superficialis - Anterior border",
        "✓ Pronator teres - Upper lateral surface",
        "✓ Flexor pollicis longus - Anterior shaft",
        "✓ Oblique cord - Lower tuberosity"
    ],
    "DISTAL ATTACHMENTS",
    [
        "✓ Pronator quadratus - Anterior lower shaft",
        "✓ Extensor carpi radialis longus - Styloid groove",
        "✓ Extensor carpi radialis brevis - Styloid groove",
        "✓ Extensor pollicis longus - Dorsal tubercle",
        "✓ Brachioradialis - Styloid process (insertion)",
        "✓ Abductor pollicis longus - Styloid area"
    ]
)

# Slide 46: Ligamentous Attachments - Complete List
add_content_slide(prs, "All Ligamentous Attachments", [
    "• ANNULAR LIGAMENT: Encircles radial head and neck",
    "• QUADRATE LIGAMENT: Supports neck of radius",
    "• INTEROSSEOUS MEMBRANE: Spans between radius and ulna",
    "• OBLIQUE CORD: Upper forearm support",
    "• RADIAL COLLATERAL LIGAMENT: Elbow joint stability",
    "• LATERAL CARPAL LIGAMENTS: From styloid process",
    "• PALMAR/DORSAL RADIO-CARPAL LIGAMENTS: Wrist support",
    "• TRIANGULAR ARTICULAR DISC: Inferior radio-ulnar joint"
])

# Slide 47: Functional Anatomy - Movement
add_content_slide(prs, "Functional Anatomy - Movements & Ranges", [
    "• FLEXION at elbow:",
    "    ◆ Biceps brachii primary action",
    "    ◆ Radius head glides on capitulum",
    "",
    "• EXTENSION at elbow:",
    "    ◆ Triceps action",
    "    ◆ Radius head moves posteriorly",
    "",
    "• PRONATION:",
    "    ◆ Radius rotates medially around ulna",
    "    ◆ 90 degrees normal range",
    "",
    "• SUPINATION:",
    "    ◆ Radius rotates laterally",
    "    ◆ 90 degrees normal range"
])

# Slide 48: Functional Anatomy - Wrist Motion
add_content_slide(prs, "Functional Anatomy - Wrist Movements", [
    "• WRIST FLEXION:",
    "    ◆ Radio-carpal joint motion",
    "    ◆ ~80-90 degrees normal range",
    "",
    "• WRIST EXTENSION:",
    "    ◆ Radio-carpal joint motion",
    "    ◆ ~70-80 degrees normal range",
    "",
    "• RADIAL DEVIATION:",
    "    ◆ Toward thumb side",
    "    ◆ ~20 degrees normal range",
    "",
    "• ULNAR DEVIATION:",
    "    ◆ Toward ulna side",
    "    ◆ ~30 degrees normal range"
])

# Slide 49: Bones of the Forearm - Comparative Anatomy
add_two_column_slide(prs, "Radius vs Ulna - Comparative Features",
    "RADIUS - LATERAL BONE",
    [
        "• Lateral position in supination",
        "• Disc-like proximal head",
        "• Rotates around ulna",
        "• Styloid extends beyond ulna",
        "• Forms lateral wrist joint",
        "• Smaller proximal end",
        "• Larger distal end"
    ],
    "ULNA - MEDIAL BONE",
    [
        "• Medial position in supination",
        "• Olecranon process proximally",
        "• Stationary - radius rotates around it",
        "• Styloid does not extend as far",
        "• Excluded from wrist joint",
        "• Larger proximal end",
        "• Smaller distal end"
    ]
)

# Slide 50: Summary & Key Learning Points
add_highlight_slide(prs, "Key Summary Points - Essential Knowledge",
    "Master these points for complete radius anatomy understanding:",
    [
        "✓ Radius is LATERAL forearm bone in supination",
        "✓ Three regions: Head (disc-like), Shaft (triangular), Styloid process",
        "✓ Articulates with HUMERUS, ULNA, SCAPHOID, and LUNATE",
        "✓ Styloid process extends BEYOND ulna's styloid",
        "✓ Dorsal tubercle (Lister's) is palpable landmark",
        "✓ ULNA is EXCLUDED from wrist/radio-carpal joint",
        "✓ Critical for pronation, supination, and hand support"
    ]
)

# Save presentation
output_path = 'Radius_Anatomy_Seminar_Complete.pptx'
prs.save(output_path)
print(f"✓ Presentation created successfully!")
print(f"✓ File: {output_path}")
print(f"✓ Total slides: {len(prs.slides)}")
print(f"✓ Content: 50 comprehensive slides")
print(f"✓ Features: Professional formatting, color scheme, organized sections")
