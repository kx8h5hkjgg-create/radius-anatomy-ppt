from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(192, 0, 60)  # Pink/Maroon
SECONDARY_COLOR = RGBColor(100, 149, 237)  # Cornflower Blue
TEXT_COLOR = RGBColor(40, 40, 40)  # Dark Gray
ACCENT_COLOR = RGBColor(220, 20, 60)  # Crimson

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 250)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Title text
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(5)
    p.space_after = Pt(5)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.line_spacing = 1.3
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add two-column layout slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 250)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, item in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.5), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    for i, item in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "THE RADIUS BONE", "Comprehensive Anatomy Guide for Students")

# Slide 2: Introduction
add_content_slide(prs, "Introduction to the Radius", [
    "• The radius is the LATERAL bone of the forearm",
    "• Pre-axial bone of lower limb",
    "• Corresponds with the tibia",
    "• Forms the lateral (thumb) side of forearm",
    "• Works with ulna for pronation and supination",
    "• Key articulation with humerus (elbow joint)"
])

# Slide 3: General Overview
add_content_slide(prs, "Basic Structure", [
    "• Three main regions:",
    "    1. Upper (proximal) end",
    "    2. Shaft (diaphysis)",
    "    3. Lower (distal) end",
    "• Functions:",
    "    - Articulates with humerus at elbow",
    "    - Rotates around ulna",
    "    - Supports wrist and hand"
])

# Slide 4: Upper End - Overview
add_content_slide(prs, "Upper End (Proximal End)", [
    "• Consists of three parts:",
    "    1. Head (caput radii)",
    "    2. Neck",
    "    3. Radial tuberosity",
    "• Articulates with humerus and ulna",
    "• Supports pronation and supination"
])

# Slide 5: The Radial Head
add_content_slide(prs, "The Radial Head (Caput)", [
    "• DISC-LIKE structure with concave upper surface",
    "• Peripheral margin broader on MEDIAL side",
    "• Upper surface: Articulates with capitulum of humerus",
    "• Forms humero-radial part of elbow joint",
    "• Peripheral margin: Articulates medially with radial notch of ulna",
    "• Rest of margin: Encircled by annular ligament",
    "• Forms superior radio-ulnar joint"
])

# Slide 6: Radial Head Details
add_content_slide(prs, "Radial Head - Articular Surfaces", [
    "• Concave disc-like upper surface",
    "    - Articulates with capitulum of humerus",
    "• Peripheral margin",
    "    - Broader on medial side",
    "    - Articulates with radial notch of ulna",
    "• Annular ligament encircles the margin",
    "• Posterior surface palpable in lateral depressed area",
    "• Moves during pronation and supination of forearm"
])

# Slide 7: Neck of Radius
add_content_slide(prs, "The Neck (Collum)", [
    "• Constricted area below the head",
    "• Encircled by lower part of annular ligament",
    "• Separated by synovial protrusion of superior radio-ulnar joint",
    "• Supported by QUADRATE LIGAMENT",
    "• Derived from interlacement of distal border of annular ligament",
    "• Provides attachment points for supporting structures"
])

# Slide 8: Radial Tuberosity
add_content_slide(prs, "Radial Tuberosity", [
    "• Located on MEDIAL side of lower part of neck",
    "• TWO distinct parts:",
    "    1. POSTERIOR part: Rough surface",
    "       - Receives insertion of TENDON OF BICEPS BRACHII",
    "       - Twisted manner of insertion",
    "    2. ANTERIOR part: Smooth surface",
    "       - Separated from biceps tendon by bursa",
    "       - Bursa for supination of forearm",
    "• Lower end: Attachment to OBLIQUE CORD"
])

# Slide 9: Shaft Overview
add_content_slide(prs, "The Shaft (Diaphysis)", [
    "• TRIANGULAR on cross-section in middle third",
    "• THREE BORDERS:",
    "    1. Anterior border",
    "    2. Posterior border",
    "    3. Interosseous (medial) border",
    "• THREE SURFACES:",
    "    1. Anterior surface",
    "    2. Posterior surface",
    "    3. Lateral surface"
])

# Slide 10: Anterior Border
add_content_slide(prs, "Anterior Border", [
    "• OBLIQUE in upper part",
    "• Slopes downward and laterally from radial tuberosity",
    "• Lower part: VERTICAL",
    "• Extends as prominent ridge in distal fourth of shaft",
    "• Continuous with anterior border of styloid process",
    "• UPPER OBLIQUE BORDER:",
    "    - Gives origin to RADIAL HEAD OF FLEXOR DIGITORUM SUPERFICIALIS",
    "• LOWER VERTICAL RIDGE:",
    "    - Attachment to lateral end of EXTENSOR RETINACULUM"
])

# Slide 11: Posterior Border
add_content_slide(prs, "Posterior Border", [
    "• PROMINENT in middle third",
    "• ILL-DEFINED in upper and lower parts",
    "• Traced above: Slopes upward and medially",
    "• Reaches: POSTERO-INFERIOR part of radial tuberosity",
    "• Provides attachment for muscular structures",
    "• Important landmark for anatomical identification"
])

# Slide 12: Interosseous Border
add_content_slide(prs, "Interosseous (Medial) Border", [
    "• SHARP border",
    "• Extends from below radial tuberosity to lower end",
    "• Connected to ulna by INTEROSSEOUS MEMBRANE",
    "• Upper margin: FREE",
    "• Gap between radius and oblique cord:",
    "    - Transmits POSTERIOR INTEROSSEOUS VESSELS",
    "    - NOT the nerves",
    "• Lower margin: Continuous with capsule ligament",
    "    - INFERIOR RADIO-ULNAR JOINT"
])

# Slide 13: Anterior Surface
add_content_slide(prs, "Anterior Surface of Shaft", [
    "• Intervenes between anterior and interosseous borders",
    "• GENTLY CONCAVE",
    "• Nutrient foramen location:",
    "    - Pointing towards elbow",
    "    - Situated near middle of shaft",
    "• Attachment points for muscles",
    "• Important for fluid dynamics during pronation"
])

# Slide 14: Posterior Surface
add_content_slide(prs, "Posterior Surface of Shaft", [
    "• Between posterior and interosseous borders",
    "• Gives origin to FLEXOR POLLICIS LONGUS",
    "• Lower one-fourth:",
    "    - Surface gives origin to flexor pollicis longus",
    "• Upper two-thirds:",
    "    - Anterior surface along with triangular medial area",
    "    - In front of interosseous membrane",
    "    - Insertion to PRONATOR QUADRATUS"
])

# Slide 15: Lateral Surface
add_content_slide(prs, "Lateral Surface of Shaft", [
    "• GENTLY CONVEX",
    "• Summit at middle of shaft",
    "• Presents rough impressions",
    "• Upper part: Encrouches impression for PRONATOR TERES insertion",
    "• Receives insertion of SUPINATOR muscle",
    "• Deep part of supinator muscle attachment",
    "• Contains deep branch of RADIAL NERVE (posterior interosseous nerve)"
])

# Slide 16: Lower End - Overview
add_content_slide(prs, "Lower End (Distal End) - Overview", [
    "• WIDEST part of the bone",
    "• Presents FOUR surfaces:",
    "    1. Lateral (Styloid process)",
    "    2. Medial (Ulnar notch)",
    "    3. Anterior",
    "    4. Posterior (with grooves)",
    "• INFERIOR CARPAL ARTICULAR SURFACE",
    "• Important for wrist and hand movement"
])

# Slide 17: Lateral Surface - Styloid Process
add_content_slide(prs, "Lateral Surface - Styloid Process", [
    "• ROUGH surface",
    "• Projects downward as STYLOID PROCESS",
    "• Extends BEYOND styloid process of ulna",
    "• Styloid process tip: Attachment to LATERAL CARPAL LIGAMENT",
    "• Proximal to styloid process:",
    "    - Receives insertion of BRACHIORADIALIS",
    "    - Crossed obliquely by ABDUCTOR POLLICIS LONGUS",
    "    - Crossed by EXTENSOR POLLICIS BREVIS tendons"
])

# Slide 18: Medial Surface - Ulnar Notch
add_content_slide(prs, "Medial Surface - Ulnar Notch", [
    "• Distal to triangular area for deep pronator quadratus fibers",
    "• Presents ULNAR NOTCH",
    "    - For articulation with HEAD OF ULNA",
    "    - Forms INFERIOR RADIO-ULNAR JOINT",
    "• Junction point:",
    "    - Between ulnar notch and carpal articular surface",
    "    - Gives attachment to BASE OF TRIANGULAR ARTICULAR DISC",
    "• Apex of disc: Fixed to depression between head of ulna and styloid process",
    "• THE ULNA IS EXCLUDED FROM WRIST JOINT FORMATION"
])

# Slide 19: Anterior Surface - Lower End
add_content_slide(prs, "Anterior Surface of Lower End", [
    "• Represented by THICK PROMINENT RIDGE",
    "• PALPABLE through overlying tendons",
    "• Attachment to PALMAR RADIO-CARPAL LIGAMENT",
    "• PULSATION OF RADIAL ARTERY:",
    "    - Felt against this surface",
    "    - Distal to PRONATOR QUADRATUS",
    "• Important clinical landmark for pulse assessment",
    "• Smooth surface for tendon gliding"
])

# Slide 20: Posterior Surface - Dorsal Tubercle
add_content_slide(prs, "Posterior Surface - Dorsal Features", [
    "• PALPABLE DORSAL TUBERCLE (of Lister)",
    "• Displays THREE LONGITUDINAL GROOVES:",
    "    1. LATERAL groove: Wide, lateral to dorsal tubercle",
    "    2. TWO MEDIAL grooves: Medial to tubercle",
    "• Lateral groove lodges:",
    "    - EXTENSOR CARPI RADIALIS LONGUS",
    "    - EXTENSOR CARPI RADIALIS BREVIS",
    "    - Ridge intervening between them"
])

# Slide 21: Posterior Surface - Medial Grooves
add_content_slide(prs, "Posterior Surface - Medial Grooves", [
    "• Groove just medial to dorsal tubercle:",
    "    - CONSPICUOUS groove",
    "    - Transmits tendon of EXTENSOR POLLICIS LONGUS",
    "    - Uses dorsal tubercle as pulley before reaching thumb",
    "• Groove more medially (shallow):",
    "    - Occupied by EXTENSOR DIGITORUM",
    "    - More deeply by EXTENSOR INDICIS",
    "    - Along with POSTERIOR INTEROSSEOUS (deep radial) NERVE"
])

# Slide 22: Posterior Surface - Extensor Retinaculum
add_content_slide(prs, "Posterior Surface - Additional Structures", [
    "• ALL structures above pass beneath EXTENSOR RETINACULUM",
    "• Gives slip of attachment to dorsal tubercle",
    "• Structures crossing this surface:",
    "    - Multiple extensor tendons",
    "    - Deep branch of radial nerve",
    "• Important for hand movement and sensation",
    "• Forms pulley system for extensor tendons"
])

# Slide 23: Inferior Carpal Articular Surface
add_content_slide(prs, "Inferior Carpal Articular Surface", [
    "• CONCAVE surface",
    "• SUBDIVIDED by ridge into:",
    "    1. LATERAL triangular area",
    "    2. MEDIAL quadrilateral area",
    "• Lateral area: Articulates with SCAPHOID",
    "• Medial area: Articulates with LUNATE",
    "• Together form the RADIO-CARPAL or WRIST JOINT",
    "• Fibrous capsule attached along periphery"
])

# Slide 24: Articular Disc
add_content_slide(prs, "Articular Disc of Inferior Radio-Ulnar Joint", [
    "• TRIANGULAR FIBROUS DISC",
    "• Attached along periphery of carpal articular surface",
    "• At ulnar notch of radius:",
    "    - Attached to anterior and posterior margins",
    "    - Attached to articular disc of inferior radio-ulnar joint",
    "• Apex fixed to depression:",
    "    - Between inferior articular surface of head of ulna",
    "    - And styloid process of ulna",
    "• Important for stability of radio-ulnar joint"
])

# Slide 25: Side Determination
add_content_slide(prs, "Side Determination - How to Identify", [
    "• DISC-LIKE HEAD OF RADIUS",
    "    - Place above",
    "    - Gentle concavity of shaft in FRONT",
    "• STYLOID PROCESS",
    "    - Projects LATERALLY",
    "• These features determine the SIDE OF THE BONE",
    "• Concave shaft surface and lateral styloid process = lateral side",
    "• Clinical importance for identification in anatomy"
])

# Slide 26: Key Articulations
add_content_slide(prs, "Key Articulations of Radius", [
    "• PROXIMAL ARTICULATIONS:",
    "    - With capitulum of humerus (humero-radial joint)",
    "    - With radial notch of ulna (superior radio-ulnar joint)",
    "• DISTAL ARTICULATIONS:",
    "    - With scaphoid and lunate (wrist/radio-carpal joint)",
    "    - With head of ulna (inferior radio-ulnar joint)",
    "• All joints essential for arm and hand function"
])

# Slide 27: Muscular Attachments Summary
add_two_column_slide(prs, "Muscular Attachments", [
    "PROXIMAL ATTACHMENTS:",
    "• Biceps brachii (tuberosity)",
    "• Supinator muscle",
    "• Flexor digitorum superficialis (anterior border)",
    "• Pronator teres (lateral surface)",
    "• Flexor pollicis longus (posterior surface)"
], [
    "DISTAL ATTACHMENTS:",
    "• Pronator quadratus (anterior surface)",
    "• Extensor carpi radialis longus",
    "• Extensor carpi radialis brevis",
    "• Extensor pollicis longus",
    "• Brachioradialis (styloid process)"
])

# Slide 28: Ligamentous Attachments
add_content_slide(prs, "Ligamentous Attachments", [
    "• ANNULAR LIGAMENT: Encircles radial head and neck",
    "• QUADRATE LIGAMENT: Supports neck",
    "• INTEROSSEOUS MEMBRANE: Connects to ulna",
    "• LATERAL CARPAL LIGAMENT: Attaches to styloid process",
    "• PALMAR RADIO-CARPAL LIGAMENT: Anterior surface",
    "• TRIANGULAR DISC: Inferior radio-ulnar joint",
    "• FIBROUS CAPSULE: Surrounds wrist joint"
])

# Slide 29: Clinical Significance
add_content_slide(prs, "Clinical Significance", [
    "• RADIAL PULSE: Palpated on anterior surface near wrist",
    "• COLLES' FRACTURE: Common distal radius fracture",
    "• SMITH'S FRACTURE: Reverse of Colles' fracture",
    "• MONTEGGIA FRACTURE: Radius fracture + ulnar dislocation",
    "• POSTERIOR INTEROSSEOUS NERVE: Can be compressed",
    "• SUPINATION/PRONATION: Compromised with fractures",
    "• WRIST DISORDERS: Affect radio-carpal joint"
])

# Slide 30: Summary & Key Points
add_content_slide(prs, "Summary - Key Points to Remember", [
    "✓ Radius is LATERAL bone of forearm",
    "✓ Three regions: Head, Shaft, Styloid process",
    "✓ Articulates with HUMERUS, ULNA, and CARPAL BONES",
    "✓ Styloid process extends BEYOND ulna styloid",
    "✓ Dorsal tubercle is important landmark",
    "✓ Ulna is EXCLUDED from wrist joint",
    "✓ Critical for pronation, supination, and hand support"
])

# Save presentation
output_path = 'Radius_Anatomy_Seminar.pptx'
prs.save(output_path)
print(f"Presentation created successfully: {output_path}")
print(f"Total slides: {len(prs.slides)}")
